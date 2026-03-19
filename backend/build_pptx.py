#!/usr/bin/env python3
"""
Gofore PPTX builder — käyttää Gofore_Template.pptx -pohjaa
Kutsutaan: python3 build_pptx.py '<json>' '<output_path>'
"""
import sys
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Gofore_Template.pptx")
# Fallback: etsi CWD:stä jos ei löydy skriptin vierestä
if not os.path.exists(TEMPLATE_PATH):
    alt = os.path.join(os.getcwd(), "Gofore_Template.pptx")
    if os.path.exists(alt):
        TEMPLATE_PATH = alt

# Gofore värit
C = {
    "deepBlue":    RGBColor(0x0C, 0x23, 0x40),
    "digitalBlue": RGBColor(0x1B, 0x6C, 0xA8),
    "orange":      RGBColor(0xE8, 0x52, 0x1A),
    "mint":        RGBColor(0x3B, 0xBF, 0xAD),
    "white":       RGBColor(0xFF, 0xFF, 0xFF),
    "grey":        RGBColor(0x8C, 0x9B, 0xAA),
    "silver":      RGBColor(0xD3, 0xD9, 0xDF),
    "light":       RGBColor(0xEE, 0xF1, 0xF3),
    "red":         RGBColor(0xC0, 0x39, 0x2B),
    "yellow":      RGBColor(0xE6, 0x7E, 0x22),
    "green":       RGBColor(0x27, 0xAE, 0x60),
}

FONT = "Cadiz"

# Layout indeksit templatessa
L = {
    "cover":     0,   # Cover slide circle only
    "bullets":   5,   # Title + content simple
    "two_col":   6,   # Title + content 2 column
    "three_col": 7,   # Title + content 3 column
    "timeline":  47,  # Headline + Timeline/table placeholder
    "section":   29,  # Section Break Title dark blue
    "end":       53,  # End slide simple
}


def set_text(tf, text, font_name=FONT, font_size=None, bold=False, color=None, align=None):
    """Aseta tekstikehyksen teksti, tyhjennä ensin."""
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color


def add_slide(prs, layout_idx):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def get_ph(slide, idx):
    """Hae placeholder indeksin mukaan."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def hide_unused_ph(slide, used_indices):
    """Poista käyttämättömät placeholderit kokonaan — estää layout-tason tekstin vuotamisen."""
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if idx not in used_indices:
            # Poista XML-elementti kokonaan slidelta
            sp = ph._element
            sp.getparent().remove(sp)


def build_title_slide(prs, d):
    """Layout 0: Cover slide circle only"""
    slide = add_slide(prs, L["cover"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("title", "Projektisuunnitelma"),
                 font_size=40, bold=True, color=C["white"])

    tagline_ph = get_ph(slide, 14)
    if tagline_ph and d.get("tagline"):
        set_text(tagline_ph.text_frame, d["tagline"], font_size=16, color=C["orange"])

    meta_ph = get_ph(slide, 15)
    if meta_ph:
        meta_parts = []
        if d.get("meta"):
            meta_parts.append(d["meta"])
        if d.get("projectLead"):
            meta_parts.append("Projektin johtaja: " + d["projectLead"])
        set_text(meta_ph.text_frame, "  |  ".join(meta_parts), font_size=12, color=C["white"])

    hide_unused_ph(slide, {0, 14, 15})


def build_bullets_slide(prs, d, def_label):
    """Layout 5: Title + content simple — käytä isoa sisältöaluetta."""
    slide = add_slide(prs, L["bullets"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])

    used = {0}

    # Käytä isoa sisältöaluetta (idx=10, 4.37" korkea) bulleteille
    # idx=16 on pieni ingressi-alue (0.63") — käytä vain jos on erillinen note
    main_ph = get_ph(slide, 10)
    if main_ph and d.get("bullets"):
        used.add(10)
        tf = main_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        bullets = d["bullets"][:10]

        total_chars = sum(len(b) for b in bullets)
        if len(bullets) > 7 or total_chars > 600:
            font_size = 12
        elif len(bullets) > 5 or total_chars > 400:
            font_size = 13
        else:
            font_size = 14

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = 0
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = bullet
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = C["deepBlue"]

    # Note/ingressi → pieni alue idx=16
    if d.get("note"):
        note_ph = get_ph(slide, 16)
        if note_ph:
            set_text(note_ph.text_frame, d["note"], font_size=12, color=C["grey"])
            used.add(16)

    hide_unused_ph(slide, used)


def build_two_col_slide(prs, d, def_label):
    """Layout 6: Title + content 2 column"""
    slide = add_slide(prs, L["two_col"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])

    left = d.get("left", {})
    right = d.get("right", {})

    # Vasen kolumni: idx=17
    left_ph = get_ph(slide, 17)
    if left_ph and left:
        tf = left_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        items = left.get("items", [])[:8]
        font_size = 11 if len(items) > 6 else 13
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = left.get("title", "")
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = C["deepBlue"]
        for item in items:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(2)
            r2 = p2.add_run()
            r2.text = "• " + item
            r2.font.name = FONT
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = C["deepBlue"]

    # Oikea kolumni: idx=18
    right_ph = get_ph(slide, 18)
    if right_ph and right:
        tf = right_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        items = right.get("items", [])[:8]
        font_size = 11 if len(items) > 6 else 13
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = right.get("title", "")
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = C["deepBlue"]
        for item in items:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(2)
            r2 = p2.add_run()
            r2.text = "• " + item
            r2.font.name = FONT
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = C["deepBlue"]

    used = {0}
    if left_ph and left: used.add(17)
    if right_ph and right: used.add(18)
    hide_unused_ph(slide, used)


def build_cards_slide(prs, d, def_label):
    """Piirretään kortit käsin tasaiseen ruudukkoon — ei templateplaceholdereja."""
    from pptx.util import Inches, Pt

    slide = add_slide(prs, L["bullets"])  # Käytetään yksinkertaista layoutia
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])

    # Piilota kaikki muut placeholderit
    hide_unused_ph(slide, {0})

    cards = d.get("cards", [])[:4]  # Max 4 korttia
    if not cards:
        return

    level_colors = {
        "high": C["red"], "korkea": C["red"],
        "medium": C["yellow"], "keski": C["yellow"],
        "low": C["green"], "matala": C["green"],
    }

    nc = len(cards)
    margin = 0.4
    gap = 0.25
    total_w = 12.5 - 2 * margin
    card_w = (total_w - (nc - 1) * gap) / nc
    card_h = 4.0
    top_y = 2.1  # Otsikko-placeholder loppuu ~1.86 — riittävä marginaali

    for i, card in enumerate(cards):
        x = margin + i * (card_w + gap)
        col = level_colors.get(str(card.get("level", "")).lower(), C["digitalBlue"])

        # Kortin tausta
        bg = slide.shapes.add_shape(1, Inches(x), Inches(top_y), Inches(card_w), Inches(card_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C["white"]
        bg.line.color.rgb = C["silver"]
        bg.line.width = Pt(1)
        bg.shadow.inherit = False

        # Väripalkkki ylhäällä
        bar = slide.shapes.add_shape(1, Inches(x), Inches(top_y), Inches(card_w), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        # Ikoni + otsikko
        title_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top_y + 0.15), Inches(card_w - 0.3), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = (card.get("icon", "") + " " + card.get("title", "")).strip()
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = col

        # Kuvaus
        if card.get("desc"):
            desc_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top_y + 0.8), Inches(card_w - 0.3), Inches(card_h - 1.0))
            tf2 = desc_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = card["desc"]
            r2.font.name = FONT
            r2.font.size = Pt(12)
            r2.font.color.rgb = C["grey"]


def build_table_slide(prs, d, def_label):
    """Layout 47: Headline + Timeline/table — taulukko"""
    slide = add_slide(prs, L["timeline"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])

    # Piilota muut placeholderit — taulukko piirretään käsin
    hide_unused_ph(slide, {0})

    cols = d.get("columns", [])
    rows = d.get("rows", [])
    if not cols or not rows:
        return

    # Poista olemassa olevat taulukot/sisältö content placeholderista
    # Lisää python-pptx taulukko suoraan slideen
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    n_cols = len(cols)
    n_rows = len(rows) + 1  # +1 header
    max_rows = 12  # Max rivejä jotka mahtuvat dialle
    if n_rows > max_rows + 1:
        rows = rows[:max_rows]
        n_rows = max_rows + 1

    left   = Inches(0.4)
    top    = Inches(2.0)  # Otsikon alapuolelle — EI päälle
    width  = Inches(12.3)
    # Dynaaminen rivikorkeus
    avail_h = 5.2  # 2.0 → 7.2, jätä marginaali
    row_h = min(0.4, avail_h / n_rows)
    height = Inches(row_h * n_rows)
    font_size = 10 if n_rows <= 8 else 8 if n_rows <= 12 else 7

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    col_width = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_width

    # Header rivi
    for ci, col_name in enumerate(cols):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["deepBlue"]
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = col_name
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = C["white"]

    # Data rivit
    for ri, row in enumerate(rows):
        bg = C["light"] if ri % 2 == 0 else C["white"]
        for ci, val in enumerate(row[:n_cols]):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val or "")
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = C["deepBlue"]


def build_gantt_slide(prs, d, def_label):
    """Layout 47: Headline + Timeline — auto months/weeks, capped to fit slide."""
    from pptx.util import Inches, Pt
    import math

    slide = add_slide(prs, L["timeline"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])

    hide_unused_ph(slide, {0})

    total_weeks = d.get("totalWeeks", 8)
    frozen_week = d.get("frozenWeek")
    phases = d.get("phases", [])

    # ═══ AUTO: kuukaudet vs viikot ═══
    # Jos > 12 viikkoa TAI > 10 saraketta, käytä kuukausia
    use_months = total_weeks > 12
    if use_months:
        total_cols = math.ceil(total_weeks / 4.33)  # viikot → kuukaudet
        col_labels = [f"Kk{i+1}" for i in range(total_cols)]
    else:
        total_cols = total_weeks
        col_labels = [f"Vk{i+1}" for i in range(total_cols)]

    # ═══ RAJOITUKSET ═══
    max_phases = 15  # Ei koskaan enempää kuin mahtuu dialle
    phases = phases[:max_phases]
    n_phases = len(phases) or 1

    # ═══ KOORDINAATIT ═══
    tl = 0.4
    tt = 2.0
    pcw = 3.2 if total_cols <= 12 else 2.8  # Kapeampi nimi jos paljon sarakkeita
    hh = 0.28
    max_bottom = 6.6  # Jätä tilaa legendille (dian korkeus ~7.5)
    avail_h = max_bottom - tt - hh - 0.35  # 0.35 legendille
    max_rh = 0.40
    min_rh = 0.22
    rh = max(min_rh, min(max_rh, avail_h / n_phases))
    avail_w = 12.5 - tl - pcw
    wcw = avail_w / total_cols

    def add_rect(slide, x, y, w, h, fill_rgb, text=None, font_size=9, bold=False, text_color=None, align=PP_ALIGN.CENTER):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.width = Pt(0.5)
        if text is not None:
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = str(text)
            r.font.name = FONT
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = text_color or C["white"]
        return shape

    # Header
    add_rect(slide, tl, tt, pcw, hh, C["deepBlue"], "Vaihe", font_size=9, bold=True)
    col_font = 8 if total_cols <= 12 else 7 if total_cols <= 18 else 6
    for ci, lbl in enumerate(col_labels):
        x = tl + pcw + ci * wcw
        is_frozen = frozen_week and not use_months and (ci + 1 == frozen_week)
        fill = C["red"] if is_frozen else C["deepBlue"]
        add_rect(slide, x, tt, wcw, hh, fill, lbl, font_size=col_font, bold=True)

    # Phases
    has_critical = any(ph.get("critical") for ph in phases)
    phase_font = 8 if n_phases <= 8 else 7 if n_phases <= 12 else 6
    name_max = 35 if pcw >= 3.0 else 28

    for ri, ph in enumerate(phases):
        y = tt + hh + ri * rh
        row_bg = C["light"] if ri % 2 == 0 else C["white"]
        is_crit = ph.get("critical", False)
        name_bg = RGBColor(0xFF, 0xF2, 0xEC) if is_crit else row_bg
        name_color = C["orange"] if is_crit else C["deepBlue"]
        prefix = "⬥ " if is_crit else ""
        name = ph.get("name", "")
        if len(prefix + name) > name_max:
            name = name[:name_max - len(prefix) - 1] + "…"

        add_rect(slide, tl, y, pcw, rh, name_bg,
                 prefix + name, font_size=phase_font, bold=is_crit,
                 text_color=name_color, align=PP_ALIGN.LEFT)

        # Muunna viikot sarakkeiksi
        start = ph.get("start", 1)
        end = ph.get("end", start)
        if use_months:
            start_col = max(0, int((start - 1) / 4.33))
            end_col = min(total_cols - 1, int((end - 1) / 4.33))
        else:
            start_col = start - 1
            end_col = end - 1

        for ci in range(total_cols):
            x = tl + pcw + ci * wcw
            active = start_col <= ci <= end_col
            is_frozen = frozen_week and not use_months and (ci + 1 == frozen_week)
            if active:
                fill = C["orange"] if is_crit else C["digitalBlue"]
            elif is_frozen:
                fill = RGBColor(0xFF, 0xE0, 0xD6)
            else:
                fill = row_bg
            add_rect(slide, x, y, wcw, rh, fill)

    # ═══ SELITE — varmista että mahtuu dialle ═══
    legend_y = min(tt + hh + n_phases * rh + 0.15, max_bottom)
    legend_items = [("Normaali vaihe", C["digitalBlue"])]
    if has_critical:
        legend_items.append(("Kriittinen polku", C["orange"]))
    if frozen_week:
        legend_items.append(("Muutosjäädytys", C["red"]))

    lx = tl
    for label, color in legend_items:
        dot = slide.shapes.add_shape(1, Inches(lx), Inches(legend_y), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(lx + 0.2), Inches(legend_y - 0.02), Inches(1.6), Inches(0.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(7)
        r.font.color.rgb = C["grey"]
        lx += 2.0


def build_bar_chart_slide(prs, d, def_label):
    """Pylväskaavio — python-pptx natiivi chart-objekti."""
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    slide = add_slide(prs, L["bullets"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])
    hide_unused_ph(slide, {0})

    categories = d.get("categories", ["A", "B", "C"])
    series_list = d.get("series", [{"name": "Data", "values": [1, 2, 3]}])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        vals = s.get("values", [])
        # Varmista numerot
        clean_vals = []
        for v in vals:
            try: clean_vals.append(float(v))
            except: clean_vals.append(0)
        chart_data.add_series(s.get("name", ""), clean_vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(2.1), Inches(11.5), Inches(4.2),
        chart_data
    ).chart

    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT

    # Värit
    gofore_colors = [C["digitalBlue"], C["orange"], C["mint"], C["grey"]]
    for si, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = gofore_colors[si % len(gofore_colors)]

    # Akselit
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = FONT
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.name = FONT

    # Note — näkyvä teksti kaavion alla
    if d.get("note"):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(11.5), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        note_text = d["note"][:200]  # Max 200 merkkiä
        r.text = note_text
        r.font.name = FONT
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C["deepBlue"]


def build_pie_chart_slide(prs, d, def_label):
    """Piirakkakaavio."""
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    slide = add_slide(prs, L["bullets"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])
    hide_unused_ph(slide, {0})

    slices = d.get("slices", [{"label": "A", "value": 50}, {"label": "B", "value": 50}])

    chart_data = CategoryChartData()
    chart_data.categories = [s.get("label", "") for s in slices]
    vals = []
    for s in slices:
        try: vals.append(float(s.get("value", 0)))
        except: vals.append(0)
    chart_data.add_series("", vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(2.0), Inches(2.1), Inches(8.5), Inches(4.2),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT

    # Värit
    gofore_colors = [C["deepBlue"], C["digitalBlue"], C["orange"], C["mint"],
                     C["grey"], RGBColor(0x5B, 0xA4, 0xCF), RGBColor(0xA0, 0x56, 0x8A)]
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = gofore_colors[i % len(gofore_colors)]

    # Datalabelit
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.font.name = FONT
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False

    if d.get("note"):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(11.5), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = d["note"][:200]
        r.font.name = FONT
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C["deepBlue"]


def build_line_chart_slide(prs, d, def_label):
    """Viivakaavio."""
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    slide = add_slide(prs, L["bullets"])
    title_ph = get_ph(slide, 0)
    if title_ph:
        set_text(title_ph.text_frame, d.get("heading", def_label), font_size=28, bold=True, color=C["deepBlue"])
    hide_unused_ph(slide, {0})

    categories = d.get("categories", ["1", "2", "3"])
    series_list = d.get("series", [{"name": "Data", "values": [1, 2, 3]}])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        vals = []
        for v in s.get("values", []):
            try: vals.append(float(v))
            except: vals.append(0)
        chart_data.add_series(s.get("name", ""), vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.5), Inches(2.1), Inches(11.5), Inches(4.2),
        chart_data
    ).chart

    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT

    gofore_colors = [C["digitalBlue"], C["orange"], C["mint"]]
    for si, series in enumerate(chart.series):
        series.format.line.color.rgb = gofore_colors[si % len(gofore_colors)]
        series.format.line.width = Pt(2.5)

    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = FONT
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.name = FONT

    if d.get("note"):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(11.5), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = d["note"][:200]
        r.font.name = FONT
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C["deepBlue"]


def build_end_slide(prs):
    """Layout 53: End slide simple"""
    add_slide(prs, L["end"])


def build_pptx(slide_data, slide_structure, output_path):
    prs = Presentation(TEMPLATE_PATH)

    # Poista kaikki olemassa olevat diat templatesta
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        sldIdLst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass

    for slide_def in slide_structure:
        sid = slide_def.get("id", "")
        layout = slide_def.get("layout", "bullets")
        label = slide_def.get("label", "")
        d = slide_data.get(sid, {})

        if layout == "title":
            build_title_slide(prs, d)
        elif layout == "two-col":
            build_two_col_slide(prs, d, label)
        elif layout == "cards":
            build_cards_slide(prs, d, label)
        elif layout == "table":
            build_table_slide(prs, d, label)
        elif layout == "gantt":
            build_gantt_slide(prs, d, label)
        elif layout == "bar_chart":
            build_bar_chart_slide(prs, d, label)
        elif layout == "pie_chart":
            build_pie_chart_slide(prs, d, label)
        elif layout == "line_chart":
            build_line_chart_slide(prs, d, label)
        else:
            build_bullets_slide(prs, d, label)

    build_end_slide(prs)
    prs.save(output_path)
    print(f"OK:{output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: build_pptx.py [--stdin|--file filepath|json_string] output.pptx", file=sys.stderr)
        sys.exit(1)

    try:
        print(f"TEMPLATE: {TEMPLATE_PATH} exists={os.path.exists(TEMPLATE_PATH)}", flush=True)
        print(f"CWD: {os.getcwd()}", flush=True)
        print(f"Script dir: {os.path.dirname(os.path.abspath(__file__))}", flush=True)

        if sys.argv[1] == "--stdin":
            # Lue JSON stdinistä — luotettavin tapa
            output = sys.argv[2]
            raw = sys.stdin.read()
            print(f"STDIN: {len(raw)} bytes", flush=True)
            payload = json.loads(raw)
        elif sys.argv[1] == "--file":
            json_path = sys.argv[2]
            output = sys.argv[3]
            with open(json_path, "r", encoding="utf-8") as f:
                payload = json.load(f)
        else:
            payload = json.loads(sys.argv[1])
            output = sys.argv[2]

        build_pptx(
            payload.get("slideData", {}),
            payload.get("slideStructure", []),
            output
        )
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"ERROR:{e}", file=sys.stderr)
        sys.exit(1)